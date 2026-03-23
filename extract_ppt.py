from pptx import Presentation

def extract_ppt():
    prs = Presentation("Portfolio Creation Lecture 2.pptx")
    with open("ppt_text.txt", "w", encoding="utf-8") as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"\n--- SLIDE {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    f.write(shape.text + "\n")

if __name__ == "__main__":
    extract_ppt()
